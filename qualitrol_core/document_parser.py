"""Input Parsing Layer (first node of the process map).

Turns customer-submitted files (Project Specification, Raw Email, Circuit
Drawing / SLD) into a normalized ``ParsedDocument`` with page/segment-level
text so downstream evidence extraction can cite a location.

Supported: .pdf (pypdf), .docx (python-docx), .txt/.eml/.msg/.md (plain text),
           .xlsx/.xlsm (openpyxl), .pptx (python-pptx), .csv (stdlib).
"""

from __future__ import annotations

import datetime
import os
import re
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Optional

from . import config


@dataclass
class DocSegment:
    """A locatable chunk of text (a PDF page, or a docx block)."""

    location: str  # e.g. "page 2" or "paragraph 14" or "table 1 row 3"
    text: str


@dataclass
class ParsedDocument:
    file_name: str
    file_path: str
    doc_type: str  # Project Specification / Raw Email / Drawing / SLD / Other
    segments: list[DocSegment] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n".join(seg.text for seg in self.segments)


# --------------------------------------------------------------------------- #
# Document role inference (maps to 11_Input_Document_Index "Document Role")
# --------------------------------------------------------------------------- #
def infer_doc_type(file_name: str, text: str = "") -> str:
    name = file_name.lower()
    blob = text.lower()
    if any(k in name for k in ("sld", "gsld", "drawing", "diagram", "_dwg")):
        return "Drawing / SLD"
    if name.endswith((".eml", ".msg")) or "raw email" in name or "from:" in blob[:500]:
        return "Raw Email"
    if any(k in name for k in ("spec", "tender", "scope", "requirement")):
        return "Project Specification"
    if any(k in name for k in ("boq", "quantity", "schedule")):
        return "Equipment List / BOQ"
    if name.endswith(".pdf") and any(
        k in blob for k in ("single line diagram", "legend", "circuit breaker")
    ):
        return "Drawing / SLD"
    return "Supporting"


# --------------------------------------------------------------------------- #
# Per-format parsers
# --------------------------------------------------------------------------- #
def _parse_pdf(
    path: Path,
    *,
    page_prefilter: Optional[Callable[[str], bool]] = None,
) -> list[DocSegment]:
    """Parse a PDF into one segment per page (page number preserved in location).

    ``page_prefilter`` is an optional, memory-saving screen for *very large*
    PDFs (e.g. a 2000+ page reference standard): a callable returning ``True``
    when a page's text should be kept. When supplied AND the document exceeds
    ``QUALITROL_PARSE_PREFILTER_MIN_PAGES`` pages, only pages that pass the
    screen (plus a small neighbour radius for context) are retained, so the
    whole body never has to sit in memory or flow into every downstream stage.
    It is a pure speed/RAM lever, gated so small/medium PDFs are unaffected, and
    fail-safe: an empty selection keeps every page (never returns nothing).
    """
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:  # pragma: no cover - corrupt page
            text = ""
        if text.strip():
            pages.append((i, text))

    if page_prefilter is not None:
        try:
            min_pages = max(1, int(os.getenv("QUALITROL_PARSE_PREFILTER_MIN_PAGES", "150")))
        except ValueError:
            min_pages = 150
        try:
            radius = max(0, int(os.getenv("QUALITROL_PARSE_PREFILTER_RADIUS", "1")))
        except ValueError:
            radius = 1
        if len(pages) >= min_pages:
            keep: set[int] = set()
            for idx, (_pno, text) in enumerate(pages):
                try:
                    hit = page_prefilter(text)
                except Exception:  # pragma: no cover - defensive: keep on error
                    hit = True
                if hit:
                    lo = max(0, idx - radius)
                    hi = min(len(pages), idx + radius + 1)
                    keep.update(range(lo, hi))
            # Only prune when the screen actually reduced volume; an empty or
            # all-pages result falls back to keeping everything (recall-safe).
            if keep and len(keep) < len(pages):
                pages = [pages[j] for j in sorted(keep)]

    return [DocSegment(location=f"page {i}", text=text) for i, text in pages]


def _parse_docx(path: Path) -> list[DocSegment]:
    import docx

    document = docx.Document(str(path))
    segments: list[DocSegment] = []
    for idx, para in enumerate(document.paragraphs, start=1):
        if para.text.strip():
            segments.append(DocSegment(location=f"paragraph {idx}", text=para.text))
    for t_idx, table in enumerate(document.tables, start=1):
        for r_idx, row in enumerate(table.rows, start=1):
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                segments.append(
                    DocSegment(location=f"table {t_idx} row {r_idx}", text=line)
                )
    return segments


# A plain number (optionally with thousands separators / %); used to tell a
# column-title cell (text) from a data cell (value) during header detection.
_XLSX_NUM_RE = re.compile(r"^[-+]?\d{1,3}(?:,\d{3})+(?:\.\d+)?$|^[-+]?\d+(?:\.\d+)?$")


def _xlsx_cell(value) -> str:
    """Render a cell to a clean string (ints without .0, dates as ISO)."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "Yes" if value else "No"
    if isinstance(value, float):
        return str(int(value)) if value.is_integer() else f"{value:g}"
    if isinstance(value, (datetime.datetime, datetime.date)):
        return value.isoformat()[:10]
    return str(value).strip()


def _xlsx_is_number(text: str) -> bool:
    return bool(text) and bool(_XLSX_NUM_RE.match(text.strip().rstrip("%")))


def _xlsx_pick_header(rows: list[tuple[int, list[str]]]) -> Optional[int]:
    """Choose the most likely column-header row among the first rows.

    A header row has several non-empty, mostly non-numeric (label-like) cells.
    Returns its position (index into ``rows``) or None when no row qualifies
    (e.g. a key-value or title-only sheet).
    """
    best_pos, best_score = None, -1.0
    for pos, (_idx, cells) in enumerate(rows[:12]):
        non_empty = [c for c in cells if c]
        if len(non_empty) < 2:
            continue
        text_cells = [c for c in non_empty if not _xlsx_is_number(c)]
        if len(text_cells) < max(2, int(0.6 * len(non_empty))):
            continue
        score = len(text_cells) - 0.3 * pos  # prefer wider, earlier rows
        if score > best_score:
            best_pos, best_score = pos, score
    return best_pos


def _classify_sheet(title: str, blob: str, *, key_value: bool) -> str:
    """Human-readable description of what a worksheet contains."""
    t = f"{title} {blob}".lower()

    def has(*keys: str) -> bool:
        return any(k in t for k in keys)

    if "boq" in title.lower() or has("bill of quant", "schedule of price",
                                      "schedule of quant"):
        return "Bill of Quantities / priced schedule"
    if has("margin", "list price", "unit price", "costing", "cost (") and has(
        "qty", "unit", "item", "price"
    ):
        return "Pricing / costing sheet"
    if has("qtms", "cbm system", "transformer monitoring", "tap changer",
           "oltc") and has("bushing", "winding", "dga", "oil temp", "cooling"):
        return "Transformer condition-monitoring (QTMS / CBM) schedule"
    if has("gas zone", "gas density", "gdht", "gdt-20", "sf6", "compartment"):
        return "SF6 gas-zone / sensor schedule"
    if has("coupler", "ocu", "uhf", "partial discharge", "pd sensor"):
        return "Partial-discharge equipment schedule"
    if has("sensor", "transducer", "rtd", "probe", "hall effect"):
        return "Sensor / accessory schedule"
    if has("part #", "part no", "item name", "model", "no. units", "qty", "quantity"):
        return "Equipment / parts table"
    if key_value:
        return "Technical parameter list (key-value)"
    return "Data table"


def _parse_xlsx(path: Path) -> list[DocSegment]:
    """Parse a workbook into human-readable segments.

    Instead of dumping raw pipe-joined cells (unreadable), each sheet is
    classified and its rows are rendered with their column headers as short
    "Header: value" statements (or "Label: value" for two-column key-value
    sheets). A leading per-sheet overview names the sheet type and columns so
    both a human reviewer and the downstream extractor get real context.
    """
    from openpyxl import load_workbook

    max_rows = int(os.getenv("QUALITROL_XLSX_MAX_ROWS", "2000"))
    max_cols = int(os.getenv("QUALITROL_XLSX_MAX_COLS", "40"))
    max_render = int(os.getenv("QUALITROL_XLSX_MAX_RENDER", "1200"))

    workbook = load_workbook(str(path), data_only=True, read_only=True)
    segments: list[DocSegment] = []
    for sheet in workbook.worksheets:
        rows: list[tuple[int, list[str]]] = []
        for r_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if r_idx > max_rows:
                break
            cells = [_xlsx_cell(v) for v in row[:max_cols]]
            while cells and cells[-1] == "":
                cells.pop()
            if any(cells):
                rows.append((r_idx, cells))
        if not rows:
            continue

        title = sheet.title
        widths = [sum(1 for c in cells if c) for _, cells in rows]
        max_w = max(widths)
        # Key-value layout: at most two columns and most rows are label+value.
        two_col = sum(1 for w in widths if w == 2)
        key_value = max_w <= 2 and two_col >= max(1, int(0.5 * len(rows)))
        header_pos = None if key_value else _xlsx_pick_header(rows)
        headers = rows[header_pos][1] if header_pos is not None else []

        blob = " ".join(" ".join(c for c in cells if c) for _, cells in rows[:60])
        classification = _classify_sheet(title, blob, key_value=key_value)

        overview = f"Sheet '{title}' - {classification}; {len(rows)} non-empty row(s)."
        header_names = [h for h in headers if h]
        if header_names:
            overview += " Columns: " + ", ".join(header_names) + "."
        segments.append(
            DocSegment(location=f"sheet '{title}' overview", text=overview)
        )

        rendered = 0
        for pos, (r_idx, cells) in enumerate(rows):
            non_empty = [c for c in cells if c]
            if not non_empty:
                continue
            if rendered >= max_render:
                segments.append(DocSegment(
                    location=f"sheet '{title}' note",
                    text=f"[... {len(rows) - pos} further rows omitted for brevity ...]",
                ))
                break

            if key_value:
                text = non_empty[0] if len(non_empty) == 1 else (
                    f"{non_empty[0]}: {', '.join(non_empty[1:])}"
                )
            elif header_pos is not None:
                if pos == header_pos:
                    continue  # header already summarised in the overview
                if pos < header_pos:
                    text = "; ".join(non_empty)  # title / description lines
                else:
                    pairs = []
                    for ci, val in enumerate(cells):
                        if not val:
                            continue
                        head = headers[ci] if ci < len(headers) and headers[ci] else f"col{ci + 1}"
                        pairs.append(f"{head}: {val}")
                    text = "; ".join(pairs)
            else:
                text = "; ".join(non_empty)

            if text.strip():
                segments.append(
                    DocSegment(location=f"sheet '{title}' row {r_idx}", text=text)
                )
                rendered += 1
    workbook.close()
    return segments


def _parse_pptx(path: Path) -> list[DocSegment]:
    try:
        from pptx import Presentation
    except ImportError:
        return _parse_text_fallback(path)

    prs = Presentation(str(path))
    segments: list[DocSegment] = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            segments.append(
                DocSegment(location=f"slide {s_idx}", text="\n".join(texts))
            )
    return segments


def _parse_csv(path: Path) -> list[DocSegment]:
    import csv

    raw = path.read_text(encoding="utf-8-sig", errors="replace")
    segments: list[DocSegment] = []
    reader = csv.reader(raw.splitlines())
    for r_idx, row in enumerate(reader, start=1):
        line = " | ".join(cell.strip() for cell in row if cell.strip())
        if line:
            segments.append(DocSegment(location=f"row {r_idx}", text=line))
    return segments


def _parse_eml(path: Path) -> list[DocSegment]:
    """Parse .eml email files, extracting headers + plain-text body."""
    from email import policy
    from email.parser import BytesParser

    raw = path.read_bytes()
    message = BytesParser(policy=policy.default).parsebytes(raw)
    header_lines = [
        f"From: {message.get('from', '')}",
        f"To: {message.get('to', '')}",
        f"Subject: {message.get('subject', '')}",
        f"Date: {message.get('date', '')}",
    ]
    body_part = message.get_body(preferencelist=("plain", "html"))
    body = body_part.get_content() if body_part else ""

    segments: list[DocSegment] = []
    if header_lines:
        segments.append(
            DocSegment(location="headers", text="\n".join(h for h in header_lines if h.split(": ", 1)[-1]))
        )
    if body.strip():
        blocks = [b.strip() for b in body.split("\n\n") if b.strip()]
        for b_idx, block in enumerate(blocks, start=1):
            segments.append(DocSegment(location=f"body block {b_idx}", text=block))
    return segments


def _parse_text_fallback(path: Path) -> list[DocSegment]:
    """Read any unrecognised file as plain text."""
    try:
        raw = path.read_text(encoding="utf-8-sig", errors="replace")
    except Exception:
        return []
    return [DocSegment(location="full", text=raw.strip())] if raw.strip() else []


def _parse_text(path: Path) -> list[DocSegment]:
    raw = path.read_text(encoding="utf-8-sig", errors="replace")
    segments: list[DocSegment] = []
    # Split into reasonably sized blocks on blank lines to keep locations useful.
    blocks = [b.strip() for b in raw.split("\n\n")]
    line_no = 1
    for block in blocks:
        n_lines = block.count("\n") + 1
        if block:
            segments.append(
                DocSegment(location=f"lines {line_no}-{line_no + n_lines - 1}", text=block)
            )
        line_no += n_lines + 1
    if not segments and raw.strip():
        segments.append(DocSegment(location="full", text=raw.strip()))
    return segments


def parse_document(
    path: str | Path,
    doc_type_override: Optional[str] = None,
    *,
    page_prefilter: Optional[Callable[[str], bool]] = None,
) -> Optional[ParsedDocument]:
    """Parse a single file into a ParsedDocument, or None if unsupported/empty.

    Args:
        path: Path to the file.
        doc_type_override: When provided, skips ``infer_doc_type`` and uses this
            value directly.  Useful when the caller already knows the role of the
            file (e.g. the user explicitly dropped it into the SLD upload zone).
        page_prefilter: Optional large-PDF page screen (see ``_parse_pdf``).
    """
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in config.SUPPORTED_DOC_EXTENSIONS:
        return None

    try:
        if ext == ".pdf":
            segments = _parse_pdf(path, page_prefilter=page_prefilter)
        elif ext == ".docx":
            segments = _parse_docx(path)
        elif ext in {".xlsx", ".xlsm"}:
            segments = _parse_xlsx(path)
        elif ext == ".pptx":
            segments = _parse_pptx(path)
        elif ext == ".csv":
            segments = _parse_csv(path)
        elif ext == ".eml":
            segments = _parse_eml(path)
        else:
            segments = _parse_text(path)
    except Exception as exc:  # pragma: no cover - defensive
        segments = [DocSegment(location="error", text=f"[parse error] {exc}")]

    if not segments:
        # A PDF with no extractable text is almost always a scanned / vector
        # drawing (SLD). Keep it — with an empty segment list — so the VLM
        # drawing-reader can still render and read it, instead of silently
        # dropping the file (which previously produced an empty BOQ with no
        # explanation). Other empty/unsupported files are still skipped.
        if ext == ".pdf" or doc_type_override:
            return ParsedDocument(
                file_name=path.name,
                file_path=str(path),
                doc_type=doc_type_override or infer_doc_type(path.name, ""),
                segments=[],
            )
        return None

    doc_type = doc_type_override or infer_doc_type(path.name, segments[0].text if segments else "")
    return ParsedDocument(
        file_name=path.name,
        file_path=str(path),
        doc_type=doc_type,
        segments=segments,
    )


def _normalize_line(line: str) -> str:
    return " ".join(line.split()).strip().lower()


def strip_running_boilerplate(
    docs: list[ParsedDocument],
    *,
    min_pages: int = 5,
    page_fraction: float = 0.5,
    max_line_len: int = 100,
) -> int:
    """Remove repeated running headers/footers from page-based documents.

    A line that recurs (normalized) on at least ``page_fraction`` of a document's
    pages and is short is a running header/footer (letterhead, tender number,
    document code, footer column titles), not requirement content. Dropping these
    lines cleans evidence snippets and stops the keyword net matching them, and
    it also reaches the grounded locator (both read the same segments). Mutates
    segment text in place; returns the number of distinct boilerplate lines
    removed. No-op for short documents where repetition is not conclusive.
    """
    removed = 0
    for doc in docs:
        page_segs = [
            s for s in doc.segments if s.location.lower().startswith("page ")
        ]
        n = len(page_segs)
        if n < min_pages:
            continue
        counts: Counter[str] = Counter()
        for seg in page_segs:
            for norm in {_normalize_line(ln) for ln in seg.text.splitlines()}:
                if norm:
                    counts[norm] += 1
        threshold = max(4, int(n * page_fraction))
        boiler = {
            line
            for line, c in counts.items()
            if c >= threshold and len(line) <= max_line_len
        }
        if not boiler:
            continue
        for seg in page_segs:
            seg.text = "\n".join(
                ln for ln in seg.text.splitlines()
                if _normalize_line(ln) not in boiler
            )
        removed += len(boiler)
    return removed


def parse_project_folder(
    folder: str | Path,
    sld_filenames: set[str] | None = None,
    *,
    page_prefilter: Optional[Callable[[str], bool]] = None,
) -> list[ParsedDocument]:
    """Parse every supported document in a customer submission folder.

    Args:
        folder: Path to the folder containing uploaded files.
        sld_filenames: Optional set of bare filenames (no path) that should be
            forced to ``"Drawing / SLD"`` doc_type, overriding ``infer_doc_type``.
            Populated from the SLD upload zone on the frontend.
        page_prefilter: Optional large-PDF page screen (see ``_parse_pdf``). Only
            engages for PDFs past the page threshold; small/medium files and
            non-PDF formats are unaffected.
    """
    folder = Path(folder)
    sld_set = {n.lower() for n in (sld_filenames or set())}
    docs: list[ParsedDocument] = []
    for file in sorted(folder.rglob("*")):
        if file.is_file() and file.suffix.lower() in config.SUPPORTED_DOC_EXTENSIONS:
            override = "Drawing / SLD" if file.name.lower() in sld_set else None
            # SLD/drawing PDFs are read by the VLM path, not text — never screen
            # their pages (would drop the drawing) so only prefilter other PDFs.
            pf = None if override == "Drawing / SLD" else page_prefilter
            parsed = parse_document(file, doc_type_override=override, page_prefilter=pf)
            if parsed:
                docs.append(parsed)
    return docs
